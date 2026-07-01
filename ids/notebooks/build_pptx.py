#!/usr/bin/env python3
"""ST-GAT IDS Graduation Project Presentation Generator — E-JUST 2026"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ───────────────────────────────────────────────────────────────────
BLUE      = RGBColor(0,   51,  102)
GOLD      = RGBColor(204, 153,   0)
L_BLUE    = RGBColor(214, 229, 244)
MED_BLUE  = RGBColor(31,  114, 199)
WHITE     = RGBColor(255, 255, 255)
OFF_W     = RGBColor(245, 248, 252)
GRAY      = RGBColor(140, 140, 140)
DK        = RGBColor(40,  40,  40)
GREEN     = RGBColor(0,   128,   0)
L_GREEN   = RGBColor(198, 239, 206)
DK_GREEN  = RGBColor(0,   80,    0)
ORANGE    = RGBColor(210, 100,   0)
L_ORANGE  = RGBColor(255, 235, 156)
DK_ORANGE = RGBColor(150,  60,   0)
RED       = RGBColor(192,   0,   0)
L_RED     = RGBColor(255, 199, 206)
PURPLE    = RGBColor(80,  0,   120)
L_PURPLE  = RGBColor(220, 200, 240)

# ── Paths ─────────────────────────────────────────────────────────────────────
HERE = os.path.dirname(os.path.abspath(__file__))
IMG  = os.path.join(HERE, "stgat_training_curves.png")
OUT  = os.path.join(HERE, "presentation_stgat.pptx")

# ── Presentation setup ────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)
HEADER_H = Inches(1.1)
FOOTER_Y = Inches(7.18)
FOOTER_H = Inches(0.32)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

TOTAL = 15

# ═════════════════════════════════════════════════════════════════════════════
# Low-level helpers
# ═════════════════════════════════════════════════════════════════════════════
def new_slide():
    return prs.slides.add_slide(BLANK)

def box(s, x, y, w, h, fill=WHITE, line_col=None, lw=Pt(1.5)):
    sh = s.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_col:
        sh.line.color.rgb = line_col
        sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def rbox(s, x, y, w, h, fill=WHITE, line_col=None, lw=Pt(1.5)):
    """Rounded rectangle (shape type 5)."""
    sh = s.shapes.add_shape(5, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_col:
        sh.line.color.rgb = line_col
        sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def img(s, path, x, y, w, h):
    s.shapes.add_picture(path, x, y, w, h)

# ── Text helpers ──────────────────────────────────────────────────────────────
def _tf_para(tf, text, size, bold, color, align, italic, idx=0):
    p = tf.paragraphs[idx] if idx == 0 else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p

def txb(s, text, x, y, w, h, size=18, bold=False, color=DK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    _tf_para(tf, text, size, bold, color, align, italic)
    return tb

def filled_box_txt(s, text, x, y, w, h, size=15, bold=False, txt_col=DK,
                   fill=L_BLUE, line_col=BLUE, lw=Pt(1.5),
                   align=PP_ALIGN.LEFT, italic=False):
    sh = box(s, x, y, w, h, fill=fill, line_col=line_col, lw=lw)
    tf = sh.text_frame
    tf.word_wrap = True
    _tf_para(tf, text, size, bold, txt_col, align, italic)
    return sh

def rbfill(s, text, x, y, w, h, size=15, bold=False, txt_col=DK,
           fill=L_BLUE, line_col=BLUE, lw=Pt(1.5), align=PP_ALIGN.LEFT):
    sh = rbox(s, x, y, w, h, fill=fill, line_col=line_col, lw=lw)
    tf = sh.text_frame
    tf.word_wrap = True
    _tf_para(tf, text, size, bold, txt_col, align)
    return sh

def multi_line(sh, lines, sizes=None, bolds=None, colors=None, aligns=None):
    tf = sh.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        sz  = (sizes[i]  if sizes  else 15)
        bd  = (bolds[i]  if bolds  else False)
        col = (colors[i] if colors else DK)
        al  = (aligns[i] if aligns else PP_ALIGN.LEFT)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = al
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = col
    return sh

# ── Arrow (thin vertical/horizontal rectangle) ────────────────────────────────
def arrow_v(s, cx, y1, y2, color=BLUE, w=Pt(4)):
    bw = Inches(0.04)
    b = box(s, cx - bw/2, y1, bw, y2 - y1, fill=color)
    ah = Inches(0.12)
    aw = Inches(0.12)
    tri = s.shapes.add_shape(5, cx - aw/2, y2 - ah/2, aw, ah)
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    return b

def arrow_h(s, x1, cy, x2, color=BLUE):
    bh = Inches(0.04)
    b = box(s, x1, cy - bh/2, x2 - x1, bh, fill=color)
    ah = Inches(0.12)
    tri = s.shapes.add_shape(5, x2 - ah/2, cy - ah/2, ah, ah)
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    return b

# ── Standard header / footer ─────────────────────────────────────────────────
def header(s, title, slide_n, subtitle=None):
    box(s, 0, 0, W, HEADER_H, fill=BLUE)
    acc = s.shapes.add_shape(1, 0, HEADER_H - Pt(4), W, Pt(4))
    acc.fill.solid(); acc.fill.fore_color.rgb = GOLD
    acc.line.fill.background()
    txb(s, title, Inches(0.3), Inches(0.07), Inches(12.5), Inches(0.65),
        size=25, bold=True, color=WHITE)
    if subtitle:
        txb(s, subtitle, Inches(0.3), Inches(0.72), Inches(12.0), Inches(0.32),
            size=13, color=GOLD)
    nbadge = box(s, Inches(12.75), Inches(0.3), Inches(0.42), Inches(0.42), fill=GOLD)
    txb(s, str(slide_n), Inches(12.75), Inches(0.3), Inches(0.42), Inches(0.42),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 0, FOOTER_Y, W, FOOTER_H, fill=BLUE)
    txb(s, "E-JUST Graduation Project 2026  ●  AI/IDS: ST-GAT for V2X  ●  Mostafa Mohamed Atef",
        Inches(0.3), FOOTER_Y + Pt(2), W - Inches(0.3), FOOTER_H,
        size=9, color=L_BLUE, align=PP_ALIGN.CENTER)

def col_header(s, text, x, y, w, h=Inches(0.45), fill=BLUE, tc=WHITE, sz=14):
    sh = box(s, x, y, w, h, fill=fill)
    txb(s, text, x, y, w, h, size=sz, bold=True, color=tc, align=PP_ALIGN.CENTER)
    return sh

def section_label(s, text, x=Inches(0.3), y=Inches(1.2), w=Inches(2.8)):
    sh = box(s, x, y, w, Inches(0.35), fill=GOLD)
    txb(s, text, x, y, w, Inches(0.35),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def table(s, headers, rows, x, y, w, h,
          header_fill=BLUE, header_tc=WHITE,
          row_fills=None, alt_fill=OFF_W, sz=13, h_sz=13):
    tbl = s.shapes.add_table(len(rows)+1, len(headers), x, y, w, h).table
    col_w = w // len(headers)
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(hdr)
        r.font.size = Pt(h_sz)
        r.font.bold = True
        r.font.color.rgb = header_tc
    for ri, row in enumerate(rows):
        fill = (row_fills[ri] if row_fills and ri < len(row_fills) else
                (alt_fill if ri % 2 == 0 else WHITE))
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(sz)
            r.font.color.rgb = DK
    return tbl

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
def slide_title():
    s = new_slide()
    box(s, 0, 0, W, H, fill=BLUE)

    box(s, 0, 0, Inches(0.18), H, fill=GOLD)
    box(s, 0, H - Inches(0.18), W, Inches(0.18), fill=GOLD)

    title_sh = box(s, Inches(0.5), Inches(0.9), Inches(11.5), Inches(2.8), fill=BLUE)
    multi_line(title_sh,
               ["AI-Driven Behavioural Anomaly Detection",
                "for V2X Communications",
                " ",
                "Spatio-Temporal Graph Attention Transformer",
                "(ST-GAT) with Targeted Fine-Tuning"],
               sizes=[34, 34, 6, 26, 26],
               bolds=[True, True, False, True, True],
               colors=[WHITE, WHITE, WHITE, GOLD, GOLD],
               aligns=[PP_ALIGN.LEFT]*5)

    dv = s.shapes.add_shape(1, Inches(0.5), Inches(3.85), Inches(11.0), Pt(3))
    dv.fill.solid(); dv.fill.fore_color.rgb = GOLD; dv.line.fill.background()

    info = box(s, Inches(0.5), Inches(4.0), Inches(9.0), Inches(2.4), fill=BLUE)
    multi_line(info,
               ["Mostafa Mohamed Atef",
                "AI / Machine Learning Engineer",
                " ",
                "Egypt-Japan University of Science and Technology (E-JUST)",
                "Faculty of Computer Science and Information Technology",
                "Graduation Project — 2026"],
               sizes=[28, 18, 6, 16, 16, 14],
               bolds=[True, False, False, True, False, False],
               colors=[WHITE, GOLD, WHITE, L_BLUE, L_BLUE, GRAY])

    badge = rbox(s, Inches(10.4), Inches(4.2), Inches(2.55), Inches(1.8),
                 fill=GOLD, line_col=WHITE, lw=Pt(2))
    multi_line(badge,
               ["AI / IDS", "Component"],
               sizes=[22, 18],
               bolds=[True, True],
               colors=[BLUE, BLUE],
               aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER])

    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═════════════════════════════════════════════════════════════════════════════
def slide_agenda():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Presentation Outline", 2)

    sections = [
        ("1", "Motivation",           "Why AI is needed alongside post-quantum cryptography"),
        ("2", "Dataset & Features",   "VeReMi NextGen (619,097 BSMs) + 15-feature engineering"),
        ("3", "ST-GAT Architecture",  "4-stream model: Transformer + iTransformer + GAT + MLP"),
        ("4", "Training Strategy",    "Focal Loss + vehicle-stratified split + targeted fine-tune"),
        ("5", "Results",              "F1 = 0.9913, MCC = 0.8699, AUC = 0.9923, 14/14 attacks"),
        ("6", "Framework Integration","ST-GAT as behavioural layer of the PQA-SCMS framework"),
    ]

    y0 = Inches(1.3)
    dh = Inches(0.88)
    colours = [BLUE, MED_BLUE, BLUE, MED_BLUE, BLUE, MED_BLUE]

    for i, (num, title, desc) in enumerate(sections):
        y = y0 + i * dh
        nbadge = box(s, Inches(0.3), y, Inches(0.55), Inches(0.65), fill=colours[i])
        txb(s, num, Inches(0.3), y, Inches(0.55), Inches(0.65),
            size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        content = box(s, Inches(1.05), y, Inches(11.8), Inches(0.65), fill=L_BLUE)
        multi_line(content,
                   [title, desc],
                   sizes=[18, 13],
                   bolds=[True, False],
                   colors=[BLUE, DK])
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MOTIVATION
# ═════════════════════════════════════════════════════════════════════════════
def slide_motivation():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Why Do We Need AI on Top of Cryptography?", 3,
           subtitle="Cryptography verifies WHO sent the message — AI verifies WHETHER the content is physically plausible")

    col_header(s, "PQA-SCMS Layer  (Team's Cryptographic Work)", Inches(0.3), Inches(1.25), Inches(6.1))
    left = box(s, Inches(0.3), Inches(1.7), Inches(6.1), Inches(3.0), fill=L_BLUE, line_col=BLUE)
    multi_line(left,
               ["✔  Dilithium-2 for long-lived certificates",
                "✔  ECDSA retained for high-frequency CAMs",
                "✔  Kyber-768 for session key exchange",
                "✔  100% replay & signature-forgery detection",
                "✔  97.4% overhead reduction vs full PQC",
                " ",
                "⚠  Cannot detect content falsification by a vehicle",
                "    with a VALID post-quantum credential"],
               sizes=[15, 15, 15, 15, 15, 6, 15, 15],
               bolds=[False]*6 + [True, True],
               colors=[DK_GREEN, DK_GREEN, DK_GREEN, DK_GREEN, DK_GREEN,
                       WHITE, RED, RED])

    col_header(s, "ST-GAT AI-IDS Layer  (This Work)", Inches(6.75), Inches(1.25), Inches(6.25), fill=RED)
    right = box(s, Inches(6.75), Inches(1.7), Inches(6.25), Inches(3.0),
                fill=RGBColor(255, 240, 240), line_col=RED)
    multi_line(right,
               ["✔  Analyses behavioural patterns in BSM content",
                "✔  Detects 14 attack types simultaneously",
                "✔  Position / speed / kinematic falsification",
                "✔  Sybil coordination attacks",
                "✔  Timing manipulation",
                "✔  DoS flooding",
                "✔  23.4 ms inference — real-time compatible"],
               sizes=[15]*7,
               bolds=[False]*7,
               colors=[DK]*7)

    ins = rbox(s, Inches(0.3), Inches(5.0), Inches(12.7), Inches(0.7),
               fill=RGBColor(255, 248, 220), line_col=GOLD, lw=Pt(2.5))
    multi_line(ins,
               ["Key Insight:  Cryptography  →  proves WHO sent it     "
                "AI-IDS  →  proves the content is physically possible"],
               sizes=[17],
               bolds=[True],
               colors=[BLUE],
               aligns=[PP_ALIGN.CENTER])
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ATTACK TAXONOMY
# ═════════════════════════════════════════════════════════════════════════════
def slide_attacks():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "V2X Attack Landscape — 14 Attack Types", 4,
           subtitle="Each category requires a different signal modality → motivates the 4-stream architecture")

    categories = [
        ("Position\nFalsification", BLUE,
         ["constantPositionOffset", "randomPositionOffset", "positionMirroring"],
         "Kinematic cross-check"),
        ("Speed\nFalsification", MED_BLUE,
         ["randomSpeedOffset", "zeroSpeedReport", "suddenConstantSpeed", "suddenStop"],
         "Speed vs displacement"),
        ("Kinematic\nInconsistency", RGBColor(0, 80, 140),
         ["feignedBraking", "accelerationMultiplication", "reversedHeading"],
         "Multi-sensor consistency"),
        ("Timing\nManipulation", ORANGE,
         ["timeDelayAttack", "dataReplay"],
         "Inter-message gap pattern"),
        ("Coordination\n& DoS", RED,
         ["trafficCongestionSybil", "dosAttack"],
         "Multi-vehicle spatial context"),
    ]

    col_w = Inches(2.45)
    gap   = Inches(0.12)
    x0    = Inches(0.3)
    y0    = Inches(1.25)

    for i, (cat, col, attacks, signal) in enumerate(categories):
        x = x0 + i * (col_w + gap)
        ch = box(s, x, y0, col_w, Inches(0.7), fill=col)
        txb(s, cat, x, y0, col_w, Inches(0.7),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        atk_box = box(s, x, y0 + Inches(0.7), col_w, Inches(3.2),
                      fill=L_BLUE, line_col=col, lw=Pt(1.5))
        lines = ["• " + a for a in attacks]
        multi_line(atk_box, lines, sizes=[13]*len(lines), bolds=[False]*len(lines),
                   colors=[DK]*len(lines))
        sig_box = box(s, x, y0 + Inches(3.9), col_w, Inches(0.6),
                      fill=col, line_col=None)
        txb(s, f"Signal: {signal}", x, y0 + Inches(3.9), col_w, Inches(0.6),
            size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    note = rbox(s, Inches(0.3), Inches(5.8), Inches(12.7), Inches(0.55),
                fill=RGBColor(255, 248, 220), line_col=GOLD, lw=Pt(2))
    txb(s, "⚠  No single feature representation covers all 14 categories — this is why ST-GAT uses 4 parallel specialised streams",
        Inches(0.3), Inches(5.8), Inches(12.7), Inches(0.55),
        size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DATASET
# ═════════════════════════════════════════════════════════════════════════════
def slide_dataset():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Dataset: VeReMi NextGen", 5,
           subtitle="SUMO-simulated highway scenario — 619,097 BSMs across 14 attack types + normal")

    col_header(s, "Dataset & Split Statistics", Inches(0.3), Inches(1.25), Inches(5.9))
    stats = [
        ("Total BSMs",        "619,097"),
        ("Attack types",      "14 + normal"),
        ("Scenario",          "Highway (SUMO)"),
        ("Area",              "624 m × 2,560 m"),
        ("Spatial cells",     "23 unique (300 m grid)"),
        ("Window size W",     "20 messages"),
        ("Stride S",          "3 (75% overlap)"),
        ("Split type",        "3-way vehicle-ID-stratified"),
        ("Train senders",     "~303  |  ~122,047 windows"),
        ("Val senders",       "~34   |  ~13,554 windows"),
        ("Test senders",      "~85   |  ~33,925 windows"),
        ("Normal upsampling", "TARGET_NORMAL_FRAC = 0.25"),
        ("Attack ratio",      "93.6% of windows"),
    ]
    tbl = table(s,
                ["Property", "Value"],
                stats,
                Inches(0.3), Inches(1.7), Inches(5.9), Inches(5.1),
                sz=12, h_sz=13)

    col_header(s, "Class Imbalance Challenge", Inches(6.55), Inches(1.25), Inches(6.45), fill=RED)
    ibox = box(s, Inches(6.55), Inches(1.7), Inches(6.45), Inches(2.4),
               fill=RGBColor(255, 240, 240), line_col=RED)
    multi_line(ibox,
               ['93.6% of all windows are labelled "attack"',
                " ",
                'A trivial "always predict attack" classifier:',
                "  Accuracy  = 93.6%  ← looks great",
                "  F1 Score  = 0.967  ← looks great",
                "  MCC       = 0.000  ← useless!"],
               sizes=[14, 6, 14, 14, 14, 14],
               bolds=[True, False, True, False, False, True],
               colors=[RED, WHITE, DK, DK, DK, RED])

    col_header(s, "Our Solution", Inches(6.55), Inches(4.15), Inches(6.45), fill=GREEN)
    sol = box(s, Inches(6.55), Inches(4.6), Inches(6.45), Inches(2.0),
              fill=L_GREEN, line_col=GREEN)
    multi_line(sol,
               ["Primary reliability metric: Matthews Correlation",
                "Coefficient (MCC) — immune to class imbalance",
                " ",
                "Our ST-GAT achieves  MCC = 0.8699",
                "(compared to MCC = 0 for trivial classifier)",
                " ",
                "Normal upsampling: reduces effective weight",
                "from 14.6× to ~3× — stabilises gradients"],
               sizes=[14, 14, 6, 16, 13, 6, 13, 13],
               bolds=[False, False, False, True, False, False, False, False],
               colors=[DK, DK, WHITE, DK_GREEN, DK, WHITE, DK, DK])
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FEATURE ENGINEERING
# ═════════════════════════════════════════════════════════════════════════════
def slide_features():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Feature Engineering: 10 Raw + 5 Physics + 8 Aggregates = Rich Input", 6)

    col_header(s, "10 Raw BSM Fields", Inches(0.3), Inches(1.25), Inches(4.1))
    raw_rows = [
        ("x, y",       "UTM position (m)"),
        ("v",          "Speed (m/s)"),
        ("φ",          "Heading (rad)"),
        ("a",          "Acceleration (m/s²)"),
        ("Δt",         "Inter-message gap (s)"),
        ("η_x",        "Position noise"),
        ("η_v",        "Speed noise"),
        ("η_φ",        "Heading noise"),
        ("η_a",        "Accel noise"),
    ]
    table(s, ["Field", "Description"], raw_rows,
          Inches(0.3), Inches(1.7), Inches(4.1), Inches(4.5), sz=12, h_sz=13)

    col_header(s, "+ 5 Physics-Based Features", Inches(4.65), Inches(1.25), Inches(4.4), fill=MED_BLUE)
    eng_rows = [
        ("Kinematic error",     "Position drift vs. kinematic model"),
        ("Speed consistency",   "Pos-implied speed vs. reported"),
        ("Heading consistency", "Movement dir vs. reported heading"),
        ("Accel consistency",   "Speed-delta accel vs. reported"),
        ("Spatial density",     "Unique vehicles in 300m × 300m / 0.5s"),
    ]
    table(s, ["Feature", "Detects"], eng_rows,
          Inches(4.65), Inches(1.7), Inches(4.4), Inches(3.0),
          header_fill=MED_BLUE, sz=12, h_sz=13)

    col_header(s, "+ 8 Window Aggregates (Stream D)", Inches(9.3), Inches(1.25), Inches(3.75), fill=ORANGE)
    agg_rows = [
        ("max(Δt)",         "Gap max → timeDelay"),
        ("σ(Δt)",           "Gap std → irregularity"),
        ("μ(Δt)",           "Gap mean → rate"),
        ("max(kin_err)",    "Kinematic spike"),
        ("max(pos_noise)",  "Position noise peak"),
        ("max(accel_err)",  "Accel inconsistency"),
        ("max(head_err)",   "Heading inconsistency"),
        ("max(density)",    "Sybil density peak"),
    ]
    table(s, ["Aggregate", "Attack Signal"], agg_rows,
          Inches(9.3), Inches(1.7), Inches(3.75), Inches(4.0),
          header_fill=ORANGE, sz=11, h_sz=12)

    note = rbox(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.5),
                fill=RGBColor(255, 248, 220), line_col=GOLD, lw=Pt(2))
    txb(s, "All inputs are z-scored (StandardScaler fitted on training data only — no val/test leakage). "
           "Window aggregates recomputed from scaled sequences.",
        Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.5),
        size=12, bold=False, color=BLUE, align=PP_ALIGN.CENTER, italic=True)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ARCHITECTURE OVERVIEW
# ═════════════════════════════════════════════════════════════════════════════
def slide_arch():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "ST-GAT Architecture — 4-Stream Parallel Design", 7,
           subtitle="Each stream specialises for one signal modality; outputs are fused into a single 384-d vector")

    streams = [
        ("A", BLUE,        L_BLUE,                "Sequence\n(20×15)",
         "3× Transformer\nEncoder Blocks\n\n4-head MHA\nover time steps\n\nCaptures temporal\ndynamics",
         "e_A ∈ ℝ¹²⁸"),
        ("B", MED_BLUE,    RGBColor(210,230,255), "Last BSM\n(15 features)",
         "iTransformer\n\nMHA over\nfeature dims\n\nCaptures feature\ncross-correlations",
         "e_B ∈ ℝ⁶⁴"),
        ("C", RGBColor(0,100,80), RGBColor(195,235,215), "Neighbours\n(8×15)",
         "2-Layer Graph\nAttention Network\n\nLayer 1: q = e_A\nLayer 2: q = h_1\n\nSpatial coordination",
         "e_C ∈ ℝ¹²⁸"),
        ("D", ORANGE,      L_ORANGE,              "Window\nAggregates (8)",
         "Window Aggregate\nMLP\n\nDense(64,GELU)\n× 2\n\nDirect attack\nfingerprints",
         "e_D ∈ ℝ⁶⁴"),
    ]

    cw  = Inches(2.9)
    gap = Inches(0.2)
    x0  = Inches(0.55)
    ih  = Inches(0.55)
    iy  = Inches(1.3)
    sy  = Inches(1.95)
    sh  = Inches(3.2)
    ey  = Inches(5.25)
    eh  = Inches(0.55)

    for i, (lbl, col, lcol, inp, detail, emb) in enumerate(streams):
        x = x0 + i * (cw + gap)
        ib = rbox(s, x, iy, cw, ih, fill=lcol, line_col=col, lw=Pt(1.5))
        txb(s, inp, x, iy, cw, ih, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        sh_b = box(s, x, sy, cw, Inches(0.45), fill=col)
        txb(s, f"  Stream {lbl}", x, sy, cw, Inches(0.45),
            size=16, bold=True, color=WHITE)
        det = box(s, x, sy + Inches(0.45), cw, sh - Inches(0.45), fill=lcol, line_col=col)
        txb(s, detail, x, sy + Inches(0.45), cw, sh - Inches(0.45),
            size=12, color=DK, align=PP_ALIGN.CENTER)
        emb_b = rbox(s, x, ey, cw, eh, fill=col, line_col=None)
        txb(s, emb, x, ey, cw, eh, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        arrow_v(s, x + cw/2, iy + ih, sy, color=col)
        arrow_v(s, x + cw/2, sy + sh, ey, color=col)

    fy = Inches(5.95)
    fus = rbox(s, Inches(0.35), fy, Inches(12.6), Inches(0.55),
               fill=GOLD, line_col=None)
    txb(s, "Concatenate:  [ e_A ‖ e_B ‖ e_C ‖ e_D ]  ∈  ℝ^(128+64+128+64)  =  ℝ³⁸⁴",
        Inches(0.35), fy, Inches(12.6), Inches(0.55),
        size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    oy = Inches(6.6)
    out = rbox(s, Inches(4.0), oy, Inches(5.3), Inches(0.52),
               fill=RED, line_col=None)
    txb(s, "Dense(128) → Dense(64) → Dense(1, σ) → anomaly score  p̂ ∈ [0, 1]",
        Inches(4.0), oy, Inches(5.3), Inches(0.52),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — STREAM DETAILS (A + B)
# ═════════════════════════════════════════════════════════════════════════════
def slide_streams_ab():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Streams A & B — Temporal and Feature Attention", 8)

    col_header(s, "Stream A — Temporal Transformer", Inches(0.3), Inches(1.25), Inches(6.1))
    a_arch = box(s, Inches(0.3), Inches(1.7), Inches(6.1), Inches(2.6), fill=L_BLUE, line_col=BLUE)
    multi_line(a_arch,
               ["Input:  X ∈ ℝ²⁰ˣ¹⁵  (20-message per-vehicle sequence)",
                "Linear projection → ℝ²⁰ˣ¹²⁸  +  LayerNorm",
                "3 × Transformer Encoder Block:",
                "  • 4-head Multi-Head Self-Attention  (d_k = 32)",
                "  • FFN: Dense(256, GELU) → Dense(128)",
                "  • Residual + LayerNorm",
                "GlobalAveragePooling  →  e_A ∈ ℝ¹²⁸"],
               sizes=[13, 13, 13, 13, 13, 13, 13],
               bolds=[False, False, True, False, False, False, True],
               colors=[DK]*6 + [BLUE])
    col_header(s, "What It Detects:", Inches(0.3), Inches(4.35), Inches(6.1), fill=GREEN)
    a_det = box(s, Inches(0.3), Inches(4.8), Inches(6.1), Inches(1.6), fill=L_GREEN, line_col=GREEN)
    multi_line(a_det,
               ["✔  constantPositionOffset — gradual trajectory drift over 20 steps",
                "✔  dataReplay — periodic pattern repetition in the sequence",
                "✔  suddenStop / suddenConstantSpeed — velocity discontinuity",
                "✔  feignedBraking — abrupt unphysical deceleration"],
               sizes=[13]*4, bolds=[False]*4, colors=[DK_GREEN]*4)

    col_header(s, "Stream B — Feature Attention (iTransformer Style)", Inches(6.75), Inches(1.25), Inches(6.25), fill=MED_BLUE)
    b_arch = box(s, Inches(6.75), Inches(1.7), Inches(6.25), Inches(2.6),
                 fill=RGBColor(210, 230, 255), line_col=MED_BLUE)
    multi_line(b_arch,
               ["Input:  x_T ∈ ℝ¹⁵  (last message only)",
                "Treat each of the 15 feature dims as a token",
                "Project each feature → ℝ⁶⁴",
                "2-head MHA across the FEATURE axis (not time!)",
                "  • Learns which features co-vary anomalously",
                "LayerNorm + GlobalAveragePool  →  e_B ∈ ℝ⁶⁴"],
               sizes=[13]*6,
               bolds=[False, True, False, True, False, True],
               colors=[DK]*5 + [MED_BLUE])
    col_header(s, "What It Detects:", Inches(6.75), Inches(4.35), Inches(6.25), fill=GREEN)
    b_det = box(s, Inches(6.75), Inches(4.8), Inches(6.25), Inches(1.6),
                fill=L_GREEN, line_col=GREEN)
    multi_line(b_det,
               ["✔  High reported speed + tiny position displacement (contradictory)",
                "✔  Large acceleration + tiny speed-noise channel (inconsistent sensors)",
                "✔  Any joint anomaly across sensor fields invisible per-feature",
                "✔  Helps on positionMirroring, randomSpeedOffset, reversedHeading"],
               sizes=[13]*4, bolds=[False]*4, colors=[DK_GREEN]*4)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — STREAMS C & D
# ═════════════════════════════════════════════════════════════════════════════
def slide_streams_cd():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Streams C & D — Graph Attention and Attack Fingerprints", 9)

    col_header(s, "Stream C — 2-Layer Graph Attention Network (GAT)", Inches(0.3), Inches(1.25), Inches(6.1), fill=RGBColor(0,100,80))
    c_arch = box(s, Inches(0.3), Inches(1.7), Inches(6.1), Inches(2.8),
                 fill=RGBColor(195, 235, 215), line_col=RGBColor(0,100,80))
    multi_line(c_arch,
               ["Input:  N ∈ ℝ⁸ˣ¹⁵  (up to 8 co-located neighbours, 300m/0.5s)",
                "Adjacency mask a prevents zero-padded slots from attending",
                " ",
                "GAT Layer 1:  query = e_A,  keys/values = N",
                "              h₁ = LayerNorm(e_A + MHA₁(e_A, N, a))",
                " ",
                "GAT Layer 2:  query = h₁,  keys/values = N",
                "              e_C = LayerNorm(h₁ + MHA₂(h₁, N, a)) ∈ ℝ¹²⁸"],
               sizes=[13]*8, bolds=[False, False, False, True, False, False, True, False],
               colors=[DK]*8)
    col_header(s, "What It Detects:", Inches(0.3), Inches(4.55), Inches(6.1), fill=GREEN)
    c_det = box(s, Inches(0.3), Inches(5.0), Inches(6.1), Inches(1.4), fill=L_GREEN, line_col=GREEN)
    multi_line(c_det,
               ["✔  trafficCongestionSybil — multiple phantoms in same 300m cell",
                "    (invisible to per-vehicle analysis; only detectable via graph)",
                "✔  Layer 2 integrates spatial context with vehicle's own dynamics"],
               sizes=[13, 13, 13], bolds=[False]*3, colors=[DK_GREEN, DK, DK_GREEN])

    col_header(s, "Stream D — Window Aggregate MLP  (Attack Fingerprints)", Inches(6.75), Inches(1.25), Inches(6.25), fill=ORANGE)
    d_arch = box(s, Inches(6.75), Inches(1.7), Inches(6.25), Inches(2.8),
                 fill=L_ORANGE, line_col=ORANGE)
    multi_line(d_arch,
               ["Input:  w ∈ ℝ⁸  (pre-computed window statistics)",
                "  max(Δt), σ(Δt), μ(Δt)  ←  inter-message gap fingerprint",
                "  max(kin_err), max(pos_noise), max(accel_err),",
                "  max(head_err), max(spatial_density)",
                " ",
                "Dense(64, GELU) → Dropout → Dense(64, GELU) → e_D ∈ ℝ⁶⁴",
                " ",
                "Key idea: reduce timeDelayAttack detection to a THRESHOLD,",
                "not a temporal discovery task"],
               sizes=[13, 13, 13, 13, 6, 14, 6, 13, 13],
               bolds=[False, True, False, False, False, True, False, True, False],
               colors=[DK, BLUE, DK, DK, WHITE, DK, WHITE, ORANGE, ORANGE])

    abl = rbox(s, Inches(6.75), Inches(4.55), Inches(6.25), Inches(1.8),
               fill=RGBColor(255, 240, 240), line_col=RED, lw=Pt(2.5))
    multi_line(abl,
               ["Stream D Ablation — timeDelayAttack Recall:",
                " ",
                "  Without Stream D  →  Recall = 0.14   (86% miss rate)",
                "  With Stream D     →  Recall = 0.8776 (12% miss rate)",
                " ",
                "Stream D provides an 82-point recall improvement on this attack"],
               sizes=[14, 6, 15, 15, 6, 13],
               bolds=[True, False, False, True, False, True],
               colors=[RED, WHITE, RED, DK_GREEN, WHITE, BLUE])
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TRAINING STRATEGY
# ═════════════════════════════════════════════════════════════════════════════
def slide_training():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Training Strategy: Two-Stage Pipeline", 10,
           subtitle="Stage 1: Vehicle-stratified split + Weighted Focal Loss  →  Stage 2: Targeted Fine-Tuning")

    s1 = rbox(s, Inches(0.4), Inches(1.3), Inches(5.9), Inches(3.8),
              fill=L_BLUE, line_col=BLUE, lw=Pt(2))
    multi_line(s1,
               ["Stage 1  —  Weighted Focal Training",
                " ",
                "Split:      3-way vehicle-ID-stratified",
                "            (train 72% / val 8% / test 20%)",
                "Upsample:   normal to 25% of training",
                "            (effective weight ~3× vs 14.6× raw)",
                "Optimizer:  Adam  lr = 1e-3",
                "EarlyStopping:  patience = 15  (validation_data=)",
                "ReduceLROnPlateau:  factor = 0.5,  patience = 7",
                " ",
                "Focal Loss:  ℒ = −α_t(1−p_t)^γ log(p_t)",
                "  γ = 2.0,  α = 0.75 for attacks",
                " ",
                "Per-attack sample weights:",
                "  timeDelayAttack      ×25.0",
                "  trafficCongestionSybil  ×20.0"],
               sizes=[16, 6, 13, 13, 13, 13, 13, 13, 13, 6, 14, 13, 6, 14, 13, 13],
               bolds=[True, False, True, False, True, False, False, False, False,
                      False, True, False, False, True, False, False],
               colors=[BLUE, WHITE, BLUE, DK, BLUE, DK, DK, DK, DK, WHITE,
                       BLUE, DK, WHITE, BLUE, DK, DK])

    arrow_h(s, Inches(6.4), Inches(3.3), Inches(7.0), color=GOLD)
    arrow_h(s, Inches(6.05), Inches(3.3), Inches(6.45), color=GOLD)
    mid = box(s, Inches(6.05), Inches(2.9), Inches(0.95), Inches(0.8), fill=GOLD)
    txb(s, "FT\nstep", Inches(6.05), Inches(2.9), Inches(0.95), Inches(0.8),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    s2 = rbox(s, Inches(7.1), Inches(1.3), Inches(5.85), Inches(3.8),
              fill=RGBColor(255, 240, 240), line_col=RED, lw=Pt(2))
    multi_line(s2,
               ["Stage 2  —  Targeted Fine-Tune",
                " ",
                "Motivation: timeDelayAttack is primary residual",
                "hard class (confirmed by per-class BCE diagnostic)",
                " ",
                "1.  Print per-class BCE to confirm hard class",
                "2.  Boost timeDelayAttack sample weight × 5",
                "3.  Fine-tune: 30 epochs, lr = 1e-4",
                "    Focal γ = 2.0, α = 0.75, patience = 10",
                "    validation_data = (clean val set)",
                " ",
                "Lower lr prevents forgetting Stage-1 knowledge.",
                "Same focal params maintain gradient stability.",
                "Saves best to stgat_ft.keras"],
               sizes=[16, 6, 13, 13, 6, 13, 13, 13, 13, 13, 6, 13, 13, 13],
               bolds=[True, False, True, False, False, False, True, True, False,
                      False, False, False, False, True],
               colors=[RED, WHITE, BLUE, DK, WHITE, DK, DK, BLUE, BLUE,
                       DK, WHITE, DK, DK, DK_GREEN])

    thr = rbox(s, Inches(0.4), Inches(5.25), Inches(12.55), Inches(0.55),
               fill=RGBColor(255, 248, 220), line_col=GOLD, lw=Pt(2))
    txb(s, "Post-training:  Threshold θ = 0.687 selected on val set by maximising F1 over ROC curve  |  "
           "Also: balanced threshold maximising min(spec_normal, rec_Sybil, rec_timeDelay)",
        Inches(0.4), Inches(5.25), Inches(12.55), Inches(0.55),
        size=12, color=BLUE, align=PP_ALIGN.CENTER, italic=True)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — OVERALL RESULTS
# ═════════════════════════════════════════════════════════════════════════════
def slide_results():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Results: Overall Performance on VeReMi NextGen Test Set", 11,
           subtitle="~33,925 test windows (vehicle-stratified 3-way split, θ = 0.687) — all project targets met")

    metrics = [
        ("0.9913", "F1 Score",    BLUE),
        ("0.9923", "ROC-AUC",    MED_BLUE),
        ("0.8699", "MCC",        GOLD),
        ("0.9886", "Recall",     RGBColor(0,120,0)),
        ("0.9940", "Precision",  RGBColor(0,80,120)),
        ("23.4 ms","Inference",  ORANGE),
    ]

    bw = Inches(2.0)
    bh = Inches(1.5)
    gap = Inches(0.15)
    x0  = Inches(0.35)
    y0  = Inches(1.3)

    for i, (val, label, col) in enumerate(metrics):
        xi = x0 + i * (bw + gap)
        b = rbox(s, xi, y0, bw, bh, fill=col, line_col=None)
        txb(s, val, xi, y0 + Inches(0.18), bw, Inches(0.78),
            size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(s, label, xi, y0 + Inches(0.95), bw, Inches(0.5),
            size=14, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    extra = [
        ("Accuracy",              "98.38%"),
        ("Normal Specificity",    "91.25%"),
        ("Operating Threshold",   "0.687 (tuned on val set)"),
        ("Attack Types Detected", "14 / 14"),
        ("Types with F1 ≥ 0.998", "13 / 14"),
        ("V2X latency budget",    "< 100 ms  ✔"),
    ]
    table(s, ["Metric", "Value"], extra,
          Inches(0.35), Inches(3.1), Inches(5.8), Inches(2.7),
          sz=14, h_sz=14)

    mcc = rbox(s, Inches(6.5), Inches(3.1), Inches(6.45), Inches(2.7),
               fill=RGBColor(255, 248, 220), line_col=GOLD, lw=Pt(2.5))
    multi_line(mcc,
               ["Why MCC = 0.8699 Matters",
                " ",
                "On a 93.6%-attack dataset:",
                " ",
                '• "Always predict attack" baseline:',
                "  Accuracy = 93.6%  F1 = 0.967  MCC = 0.000",
                " ",
                "• ST-GAT:",
                "  Accuracy = 98.4%  F1 = 0.991  MCC = 0.870",
                " ",
                "MCC = 0.870 proves genuine skill beyond",
                "exploiting the majority class."],
               sizes=[16, 6, 14, 6, 14, 14, 6, 14, 14, 6, 13, 13],
               bolds=[True, False, False, False, False, True, False, False, True, False, True, False],
               colors=[BLUE, WHITE, DK, WHITE, RED, RED, WHITE, DK_GREEN, DK_GREEN, WHITE, BLUE, BLUE])
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — TRAINING CURVES
# ═════════════════════════════════════════════════════════════════════════════
def slide_curves():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Training Convergence — Stage 1 & Stage 2 History", 12,
           subtitle="Clean vehicle-stratified validation set provides reliable EarlyStopping signal")

    if os.path.exists(IMG):
        img(s, IMG, Inches(0.7), Inches(1.25), Inches(8.5), Inches(4.9))
    else:
        ph = box(s, Inches(0.7), Inches(1.25), Inches(8.5), Inches(4.9),
                 fill=L_BLUE, line_col=BLUE)
        txb(s, "[ stgat_training_curves.png ]", Inches(0.7), Inches(3.5), Inches(8.5), Inches(0.7),
            size=18, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    obs = [
        (BLUE,     "2×2 Layout",
         "Row 1: Stage 1 accuracy + loss\nRow 2: Stage 2 accuracy + loss"),
        (GREEN,    "Clean Validation",
         "Val set = separate vehicle senders.\nNo overlapping windows with training."),
        (ORANGE,   "Threshold",
         "θ = 0.687 tuned on val set.\nOptimises F1 across ROC curve."),
    ]
    ya = Inches(1.3)
    for col, title, body in obs:
        tb = rbox(s, Inches(9.4), ya, Inches(3.65), Inches(1.45), fill=col, line_col=None)
        multi_line(tb, [title, " ", body],
                   sizes=[15, 5, 12],
                   bolds=[True, False, False],
                   colors=[WHITE, WHITE, WHITE])
        ya += Inches(1.6)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — PER-ATTACK RESULTS
# ═════════════════════════════════════════════════════════════════════════════
def slide_per_attack():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Per-Attack-Type Detection Results", 13,
           subtitle="Threshold θ = 0.687 applied to all classes; all 14 attack types detected")

    rows = [
        ("accelerationMultiplication", "0.9995", "0.9990", "1.0000"),
        ("constantPositionOffset",     "0.9995", "0.9990", "1.0000"),
        ("dataReplay",                 "0.9995", "0.9990", "1.0000"),
        ("dosAttack",                  "0.9980", "0.9965", "1.0000"),
        ("feignedBraking",             "0.9995", "0.9990", "1.0000"),
        ("NORMAL  (specificity)",      "0.9125", "—",      "—"),
        ("positionMirroring",          "0.9995", "0.9990", "1.0000"),
        ("randomPositionOffset",       "0.9995", "0.9990", "1.0000"),
        ("randomSpeedOffset",          "0.9995", "0.9990", "1.0000"),
        ("reversedHeading",            "0.9995", "0.9990", "1.0000"),
        ("suddenConstantSpeed",        "0.9995", "0.9990", "1.0000"),
        ("suddenStop",                 "0.9995", "0.9990", "1.0000"),
        ("⚠ timeDelayAttack",          "0.9348", "0.8776", "1.0000"),
        ("trafficCongestionSybil",     "0.9995", "0.9990", "1.0000"),
        ("zeroSpeedReport",            "0.9995", "0.9990", "1.0000"),
    ]

    row_fills = []
    for r in rows:
        nm = r[0]
        if "timeDelay" in nm:
            row_fills.append(L_ORANGE)
        elif "NORMAL" in nm:
            row_fills.append(RGBColor(210, 225, 245))
        elif float(r[1]) >= 0.999:
            row_fills.append(L_GREEN)
        else:
            row_fills.append(RGBColor(230, 248, 230))

    table(s, ["Attack Type", "F1 Score", "Recall", "Precision"],
          rows,
          Inches(0.3), Inches(1.25), Inches(8.5), Inches(5.7),
          row_fills=row_fills, sz=11, h_sz=12)

    n1 = rbox(s, Inches(9.1), Inches(1.3), Inches(3.95), Inches(1.6),
              fill=L_GREEN, line_col=GREEN, lw=Pt(2))
    multi_line(n1, ["13 / 14 Attack Types\nAchieve F1 ≥ 0.998",
                    " ",
                    "Including the hardest:\n• Sybil coordination\n• dataReplay\n• dosAttack"],
               sizes=[15, 6, 13], bolds=[True, False, False],
               colors=[DK_GREEN, WHITE, DK_GREEN])

    n2 = rbox(s, Inches(9.1), Inches(3.1), Inches(3.95), Inches(1.8),
              fill=L_ORANGE, line_col=ORANGE, lw=Pt(2))
    multi_line(n2, ["timeDelayAttack  F1 = 0.9348",
                    " ",
                    "Recall: 0.8776",
                    "(was 0.14 without Stream D)",
                    "Precision: 1.000 — zero",
                    "false alarms for this class"],
               sizes=[14, 6, 13, 13, 13, 13], bolds=[True, False, False, False, False, False],
               colors=[DK_ORANGE, WHITE, DK, DK, DK, DK])

    n3 = rbox(s, Inches(9.1), Inches(5.05), Inches(3.95), Inches(1.2),
              fill=L_BLUE, line_col=BLUE, lw=Pt(2))
    multi_line(n3, ["Normal Specificity = 91.25%",
                    " ",
                    "Threshold θ = 0.687 balances\nattack recall vs false-alarm rate"],
               sizes=[14, 6, 13], bolds=[True, False, False],
               colors=[BLUE, WHITE, DK])
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — TWO-LAYER FRAMEWORK
# ═════════════════════════════════════════════════════════════════════════════
def slide_framework():
    s = new_slide()
    box(s, 0, 0, W, H, fill=OFF_W)
    header(s, "Complete Two-Layer Security Framework", 14,
           subtitle="PQA-SCMS + ST-GAT provide disjoint, complementary coverage of the full V2X threat landscape")

    l1 = rbox(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(1.4),
              fill=BLUE, line_col=None)
    multi_line(l1,
               ["Layer 1 — PQA-SCMS:  Post-Quantum Cryptographic Authentication  (Team's Work)",
                "Dilithium-2 certs  ·  ECDSA for CAMs  ·  Kyber-768 session keys  ·  "
                "100% replay & forgery detection  ·  97.4% overhead reduction vs full-PQC"],
               sizes=[18, 13], bolds=[True, False], colors=[WHITE, L_BLUE])

    dv = box(s, W/2 - Inches(0.03), Inches(2.8), Inches(0.06), Inches(0.3), fill=GOLD)

    l2 = rbox(s, Inches(0.4), Inches(3.15), Inches(12.5), Inches(1.4),
              fill=RGBColor(180, 0, 0), line_col=None)
    multi_line(l2,
               ["Layer 2 — ST-GAT IDS:  AI-Driven Behavioural Detection  (This Work)",
                "F1 = 0.9913  ·  MCC = 0.8699  ·  AUC = 0.9923  ·  14/14 attacks  ·  23.4 ms inference"],
               sizes=[18, 13], bolds=[True, False], colors=[WHITE, RGBColor(255, 200, 200)])

    cov_rows = [
        ("Key forgery / ECDSA break",       "✔  100%",          "—"),
        ("HNDL on enrollment certs",        "✔  Prevented",      "—"),
        ("Replay (invalid signature)",       "✔  100%",          "✔  F1 = 0.9995"),
        ("Content falsification — position","—  Not detected",   "✔  F1 = 0.9995"),
        ("Content falsification — speed",   "—  Not detected",   "✔  F1 = 0.9980–0.9995"),
        ("Timing manipulation",             "—  Not detected",   "⚠  F1 = 0.9348"),
        ("Sybil attacks (valid creds)",     "—  Not detected",   "✔  F1 = 0.9995"),
        ("DoS flooding (valid creds)",      "—  Not detected",   "✔  F1 = 0.9980"),
    ]
    cov_fills = [L_BLUE, L_BLUE, L_GREEN, L_GREEN, L_GREEN, L_ORANGE, L_GREEN, L_GREEN]
    table(s, ["Attack Class", "PQA-SCMS", "ST-GAT (this work)"],
          cov_rows,
          Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.2),
          row_fills=cov_fills, sz=12, h_sz=13)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — THANK YOU
# ═════════════════════════════════════════════════════════════════════════════
def slide_thankyou():
    s = new_slide()
    box(s, 0, 0, W, H, fill=BLUE)
    box(s, 0, 0, Inches(0.2), H, fill=GOLD)
    box(s, 0, H - Inches(0.2), W, Inches(0.2), fill=GOLD)

    txb(s, "Thank You", 0, Inches(0.9), W, Inches(1.1),
        size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    div = box(s, Inches(3.5), Inches(2.1), Inches(6.3), Pt(3), fill=GOLD)

    card = rbox(s, Inches(1.5), Inches(2.3), Inches(10.3), Inches(2.6),
                fill=RGBColor(0, 40, 85), line_col=GOLD, lw=Pt(2.5))
    multi_line(card,
               ["Mostafa Mohamed Atef  —  AI / Machine Learning Engineer",
                "Egypt-Japan University of Science and Technology (E-JUST)",
                " ",
                "ST-GAT IDS for V2X     F1 = 0.9913     MCC = 0.8699",
                "ROC-AUC = 0.9923       14 / 14 Attacks     23.4 ms inference",
                " ",
                "Code:  loss_increased.ipynb    |    Paper:  paper_stgat_ids.tex"],
               sizes=[20, 15, 6, 20, 17, 6, 12],
               bolds=[True, False, False, True, True, False, False],
               colors=[WHITE, L_BLUE, WHITE, GOLD, L_BLUE, WHITE, GRAY],
               aligns=[PP_ALIGN.CENTER]*7)

    txb(s, "Questions?", 0, Inches(5.1), W, Inches(0.8),
        size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# BUILD ALL SLIDES
# ═════════════════════════════════════════════════════════════════════════════
print("Building slides...")
slide_title()
slide_agenda()
slide_motivation()
slide_attacks()
slide_dataset()
slide_features()
slide_arch()
slide_streams_ab()
slide_streams_cd()
slide_training()
slide_results()
slide_curves()
slide_per_attack()
slide_framework()
slide_thankyou()

prs.save(OUT)
print(f"Saved: {OUT}")
